from pathlib import Path
from typing import Optional

from .base import BaseConverter, ConversionResult
from ..exceptions import MissingDependencyError


class PptxConverter(BaseConverter):
    """Convert PPTX presentations to Markdown using python-pptx."""

    EXTENSIONS = {".pptx"}
    MIME_TYPES = {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    }

    @classmethod
    def accepts(cls, mime_type: str, extension: Optional[str] = None) -> bool:
        if mime_type in cls.MIME_TYPES:
            return True
        if extension and extension.lower() in cls.EXTENSIONS:
            return True
        return False

    def convert(self, path: str, **kwargs) -> ConversionResult:
        path_obj = Path(path)
        if not path_obj.exists():
            return ConversionResult(
                text="",
                error="File not found",
                error_message=f"{path} does not exist",
            )

        try:
            from pptx import Presentation
        except ImportError:
            raise MissingDependencyError(
                "python-pptx is required for PPTX conversion. "
                "Install: pip install zarishnote-ingest[pptx]"
            )

        try:
            prs = Presentation(str(path_obj))
            title = path_obj.stem
            md_lines: list[str] = []

            for i, slide in enumerate(prs.slides, 1):
                md_lines.append(f"## Slide {i}")
                md_lines.append("")

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                prefix = ""
                                if paragraph.level > 0:
                                    prefix = "  " * paragraph.level + "- "
                                md_lines.append(f"{prefix}{text}")
                                md_lines.append("")

                    if shape.has_table:
                        table = shape.table
                        for row_idx, row in enumerate(table.rows):
                            cells = [
                                cell.text.strip().replace("\n", " ")
                                for cell in row.cells
                            ]
                            md_lines.append(
                                "| " + " | ".join(cells) + " |"
                            )
                            if row_idx == 0:
                                md_lines.append(
                                    "| "
                                    + " | ".join(["---"] * len(cells))
                                    + " |"
                                )
                        md_lines.append("")

                    if hasattr(shape, "image") and shape.image:
                        md_lines.append(f"*[Image: slide {i}]*")
                        md_lines.append("")

                md_lines.append("---")
                md_lines.append("")

            text = "\n".join(md_lines)

            return ConversionResult(
                text=text,
                title=title,
                metadata={
                    "source": str(path_obj),
                    "slides": len(prs.slides),
                    "slide_count": len(prs.slides),
                },
            )
        except Exception as e:
            return ConversionResult(
                text="",
                error="PPTX conversion failed",
                error_message=str(e),
            )
